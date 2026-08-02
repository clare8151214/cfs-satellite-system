#!/usr/bin/env python3
"""Generate the research progress slide deck for the cFS FreeRTOS POC."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "研究進度-cFS-FreeRTOS-POC.pptx"

NAVY = RGBColor(20, 39, 67)
TEAL = RGBColor(0, 132, 121)
GREEN = RGBColor(80, 148, 83)
AMBER = RGBColor(220, 151, 54)
RED = RGBColor(190, 72, 72)
INK = RGBColor(36, 42, 52)
MUTED = RGBColor(95, 104, 116)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(209, 216, 226)


def set_run(run, size=20, bold=False, color=INK):
    run.font.name = "Microsoft JhengHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(slide, text, x, y, w, h, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return shape


def add_title(slide, title, subtitle=None):
    add_textbox(slide, title, 0.55, 0.32, 12.1, 0.55, size=26, color=NAVY, bold=True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.98), Inches(1.25), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, subtitle, 1.95, 0.86, 10.8, 0.35, size=12, color=MUTED)


def add_footer(slide, page):
    add_textbox(slide, f"cFS FreeRTOS POC | {page}", 10.95, 7.05, 1.8, 0.2, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items, x, y, w, h, size=17, color=INK):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = "Microsoft JhengHei"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(8)
    return shape


def add_card(slide, title, body, x, y, w, h, fill=WHITE, accent=TEAL):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, title, x + 0.22, y + 0.18, w - 0.35, 0.35, size=16, color=NAVY, bold=True)
    add_textbox(slide, body, x + 0.22, y + 0.65, w - 0.35, h - 0.75, size=12.5, color=INK)
    return card


def add_box(slide, text, x, y, w, h, fill, color=WHITE, size=13, bold=True):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return shape


def add_line(slide, x1, y1, x2, y2, color=MUTED):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    return line


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    add_textbox(slide, "研究進度", 0.75, 1.25, 7.0, 0.75, size=42, color=NAVY, bold=True)
    add_textbox(slide, "cFS 衛星系統 FreeRTOS POC 與地面站串接", 0.82, 2.05, 7.5, 0.45, size=20, color=TEAL, bold=True)
    add_bullets(
        slide,
        [
            "目標：將原本 Ubuntu ARM64 衛星 VM 的 cFS 概念，驗證到 FreeRTOS/QEMU",
            "目前：cFE 可啟動 mission app，並以 GroundSystem 接收遙測與命令 ACK",
        ],
        0.85,
        3.05,
        7.3,
        1.2,
        size=18,
    )
    add_box(slide, "FreeRTOS\nQEMU Cortex-M3", 8.7, 1.35, 3.55, 1.2, NAVY, size=20)
    add_box(slide, "cFE / OSAL / PSP", 9.15, 3.05, 2.65, 0.78, TEAL, size=16)
    add_box(slide, "cFS-GroundSystem", 8.7, 4.55, 3.55, 0.95, GREEN, size=17)
    add_line(slide, 10.45, 2.55, 10.45, 3.05, TEAL)
    add_line(slide, 10.45, 3.83, 10.45, 4.55, GREEN)
    add_footer(slide, "1")

    # 2. Goal
    slide = prs.slides.add_slide(blank)
    add_title(slide, "研究目標", "從可開機，推進到可展示的衛星任務資料流")
    add_card(slide, "原系統", "Ubuntu ARM64 VM\n作為衛星端執行環境\n依賴 Linux socket / filesystem", 0.75, 1.55, 3.55, 1.85, accent=AMBER)
    add_card(slide, "老師要求", "把衛星端作業系統方向\n切到 FreeRTOS\n保留 cFS/cFE 核心概念", 4.9, 1.55, 3.55, 1.85, accent=TEAL)
    add_card(slide, "本次 POC", "QEMU mps2-an385 Cortex-M3\nFreeRTOS + cFE\nmission app + 地面站資料流", 9.05, 1.55, 3.55, 1.85, accent=GREEN)
    add_bullets(
        slide,
        [
            "驗證 cFE 在 FreeRTOS target 可進入 OPERATIONAL",
            "補上 /cf startup script 與 static mission app",
            "用 host-side bridge 先接回 cFS-GroundSystem",
            "清楚標出目前尚未 native 化的部分",
        ],
        1.0,
        4.15,
        11.2,
        1.8,
    )
    add_footer(slide, "2")

    # 3. Architecture
    slide = prs.slides.add_slide(blank)
    add_title(slide, "目前架構", "FreeRTOS 衛星端 + host bridge + cFS-GroundSystem")
    add_box(slide, "cFS-GroundSystem\nPyQt GUI", 0.75, 1.55, 2.45, 0.85, GREEN, size=14)
    add_box(slide, "Command GUI\nUDP 1234", 0.75, 3.0, 2.45, 0.72, GREEN, size=13)
    add_box(slide, "Telemetry GUI\nUDP 2234", 0.75, 4.15, 2.45, 0.72, GREEN, size=13)
    add_box(slide, "satellite-ground-bridge.py\nhost-side POC bridge", 4.55, 2.65, 3.05, 1.35, TEAL, size=14)
    add_box(slide, "QEMU ARM Cortex-M3\nmps2-an385", 9.05, 1.25, 3.1, 0.78, NAVY, size=14)
    add_box(slide, "FreeRTOS image\ncore-mps2", 9.05, 2.35, 3.1, 0.78, NAVY, size=14)
    add_box(slide, "cFE + OSAL + PSP", 9.05, 3.45, 3.1, 0.78, TEAL, size=14)
    add_box(slide, "SAT_SAMPLE_APP\nSatellite Mission HK", 9.05, 4.55, 3.1, 0.78, GREEN, size=14)
    add_line(slide, 3.2, 3.35, 4.55, 3.35, TEAL)
    add_line(slide, 7.6, 3.35, 9.05, 4.92, TEAL)
    add_line(slide, 9.05, 4.92, 7.6, 3.35, AMBER)
    add_line(slide, 4.55, 3.55, 3.2, 4.5, AMBER)
    add_textbox(slide, "命令：0x1882", 3.25, 3.02, 1.15, 0.25, size=10, color=MUTED)
    add_textbox(slide, "遙測：0x883", 3.25, 4.25, 1.15, 0.25, size=10, color=MUTED)
    add_footer(slide, "3")

    # 4. Completed work
    slide = prs.slides.add_slide(blank)
    add_title(slide, "已完成項目", "目前已經能展示的部分")
    add_card(slide, "1. FreeRTOS cFE boot", "cFE 在 QEMU Cortex-M3 上啟動\n進入 OPERATIONAL state", 0.8, 1.35, 3.75, 1.45, accent=GREEN)
    add_card(slide, "2. /cf startup script", "嵌入 /cf/cfe_es_startup.scr\n啟動 SAT_SAMPLE_APP", 4.8, 1.35, 3.75, 1.45, accent=GREEN)
    add_card(slide, "3. Mission app 行為", "SAFE -> NOMINAL mode\n輸出 uptime / payload / battery HK", 8.8, 1.35, 3.75, 1.45, accent=GREEN)
    add_card(slide, "4. Static loader 補強", "FreeRTOS loader stub 可查 static module/symbol\n解決無 dynamic loader 問題", 0.8, 3.45, 3.75, 1.45, accent=TEAL)
    add_card(slide, "5. Ground bridge", "console telemetry -> UDP 2234\nUDP 1234 command -> ACK telemetry", 4.8, 3.45, 3.75, 1.45, accent=TEAL)
    add_card(slide, "6. GroundSystem 頁面", "新增 Satellite Mission HK\n新增 Satellite Mission command entry", 8.8, 3.45, 3.75, 1.45, accent=TEAL)
    add_footer(slide, "4")

    # 5. Demo and evidence
    slide = prs.slides.add_slide(blank)
    add_title(slide, "展示流程與驗證", "可以現場 demo 的最短路徑")
    add_bullets(
        slide,
        [
            "Terminal 1：cd ~/cFS-GroundSystem && python3 GroundSystem.py",
            "Terminal 2：cd ~/cfs-freertos-satellite && ./start-satellite-freertos-poc.sh",
            "GroundSystem Telemetry：開 Satellite Mission HK",
            "GroundSystem Command：送 Satellite Mission / Mission No-Op",
        ],
        0.85,
        1.35,
        5.7,
        2.0,
        size=16,
    )
    add_card(
        slide,
        "驗證 log 摘要",
        "SAT_MISSION_HK,1,0,0,1,0,1,0,97\n"
        "[bridge] UDP telemetry -> 127.0.0.1:2234 ...\n"
        "[bridge] UDP command <- ... pkt=0x1882 cc=0 accepted\n"
        "[bridge] UDP telemetry -> ... reason=command-accepted",
        6.9,
        1.35,
        5.45,
        2.55,
        fill=RGBColor(252, 253, 255),
        accent=AMBER,
    )
    add_card(
        slide,
        "GroundSystem 顯示欄位",
        "Command Counter\nError Counter\nMission Mode：BOOT / SAFE / NOMINAL\nMission Status：OK / LOW_BAT / CMD_ERR\nUptime Seconds / Payload Samples / Battery Percent",
        0.85,
        4.25,
        11.5,
        1.55,
        accent=GREEN,
    )
    add_footer(slide, "5")

    # 6. Limitations and next steps
    slide = prs.slides.add_slide(blank)
    add_title(slide, "限制與下一步", "目前是 POC bridge，下一步要往 native flight path 推進")
    add_card(
        slide,
        "目前限制",
        "FreeRTOS OSAL socket layer 還是 stub\n"
        "命令目前由 host bridge 接收並回 ACK\n"
        "尚未把 command packet 注入 cFE Software Bus\n"
        "filesystem 仍以 embedded startup script 為主",
        0.8,
        1.35,
        5.6,
        3.1,
        accent=RED,
    )
    add_card(
        slide,
        "下一步",
        "1. 實作 native command ingest：UART RX 或 OSAL sockets\n"
        "2. 將 telemetry 改成 TO_LAB-style native UDP\n"
        "3. 擴充 mission command 與 packet definition\n"
        "4. 評估 /cf filesystem 是否完整化",
        6.9,
        1.35,
        5.6,
        3.1,
        accent=TEAL,
    )
    add_box(slide, "結論：已從「能 boot」推進到「可展示任務資料流與地面站串接」", 1.25, 5.25, 10.85, 0.75, NAVY, size=20)
    add_footer(slide, "6")

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build_deck())
